"""PPTX parser（python-pptx）—— 抽 slide 文字 + 备注

每张幻灯片：标题 + body 文字 + speaker notes，按 \n\n 串成一篇。
"""

from __future__ import annotations

import io
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from chameleon.api.knowledge.parsers import ParsedDocument


async def parse(source: bytes | str, *, name: str) -> "ParsedDocument":
    from chameleon.api.knowledge.parsers import ParsedDocument

    if isinstance(source, str):
        raise TypeError("pptx parser requires bytes input")

    try:
        from pptx import Presentation
    except ImportError as e:  # noqa: BLE001
        raise RuntimeError("python-pptx not installed; uv add python-pptx") from e

    prs = Presentation(io.BytesIO(source))
    slides_text: list[str] = []
    for idx, slide in enumerate(prs.slides, 1):
        parts: list[str] = [f"--- 幻灯片 {idx} ---"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = shape.text.strip()
                if txt:
                    parts.append(txt)
        # 演讲者备注
        try:
            notes = (
                slide.notes_slide.notes_text_frame.text
                if slide.has_notes_slide
                else ""
            )
            if notes and notes.strip():
                parts.append(f"[备注] {notes.strip()}")
        except Exception:  # noqa: BLE001
            pass
        slides_text.append("\n".join(parts))

    text = "\n\n".join(slides_text)
    metadata: dict = {"name": name, "slide_count": len(prs.slides)}
    return ParsedDocument(text=text, metadata=metadata)
